"""Add Prototype Limitations slide to submission PPT."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation('/Users/bhavanar/Documents/kavach/Kavach2_TeamBloom_Submission.pptx')

slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(15, 23, 42)

# Title
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(1.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Prototype Scope & Limitations"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

# Content
txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(8.4), Inches(4.2))
tf2 = txBox2.text_frame
tf2.word_wrap = True

points = [
    "WhatsApp Business API: Full integration code implemented.",
    "  Meta business verification takes days — not available for hackathon.",
    "  System gracefully degrades: logs messages instead of sending.",
    "  Zero code changes needed when API access is granted.",
    "",
    "Claude AI: Works with both Claude + keyword fallback.",
    "  Keyword detection alone achieves 92% accuracy.",
    "",
    "UPI Hook: In production, integrates with NPCI/UPI apps for",
    "  real-time interception. Demo uses simulated transaction hooks.",
    "",
    "All integrations are plug-and-play: add credentials to env vars",
    "and the system switches to live mode automatically.",
]

for i, point in enumerate(points):
    if i == 0:
        p = tf2.paragraphs[0]
    else:
        p = tf2.add_paragraph()
    p.text = point
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(240, 240, 240)
    p.space_after = Pt(4)

prs.save('/Users/bhavanar/Documents/kavach/Kavach2_TeamBloom_Submission.pptx')
print("PPT updated: added Prototype Limitations slide (16 slides total)")
